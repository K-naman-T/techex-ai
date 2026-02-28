import os
import json
from pptx import Presentation

BASE_DIR = r"C:\Users\naman\Downloads\OneDrive_2026-02-27\Checked Backdrops by Backdrop Team"

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text.strip())
    return " ".join([t for t in text_content if t])

def get_unique_pptx_files(base_dir):
    unique_files = {}
    for root, dirs, files in os.walk(base_dir):
        for file in files:
            if file.endswith(".pptx"):
                # Use filename as key to avoid duplicates (e.g. TE-26-037.pptx)
                if file not in unique_files:
                    unique_files[file] = os.path.join(root, file)
    return unique_files

def main():
    files_map = get_unique_pptx_files(BASE_DIR)
    results = []
    
    print(f"Found {len(files_map)} unique PPTX files.")
    
    for filename, full_path in files_map.items():
        try:
            content = extract_text_from_pptx(full_path)
            # Try to extract Stall Number from filename (e.g. TE-26-037 -> 037)
            stall_number = "Unknown"
            if "TE-26-" in filename:
                stall_number = filename.split("TE-26-")[1].replace(".pptx", "").replace("_checked", "").replace("_Checked", "").strip()
            elif "TE 26_" in filename:
                stall_number = filename.split("TE 26_")[1].replace(".pptx", "").strip()
            
            # Basic title extraction (first line or filename)
            lines = content.split('\n')
            title = lines[0] if lines else filename
            
            results.append({
                "title": title,
                "stall_number": stall_number,
                "category": "Exhibit",
                "description": content,
                "filename": filename
            })
            print(f"Processed: {filename}")
        except Exception as e:
            print(f"Error processing {filename}: {e}")

    with open("extracted_projects.json", "w", encoding="utf-8") as f:
        json.dump(results, f, indent=2, ensure_ascii=False)
    
    print(f"Extraction complete. Saved to extracted_projects.json")

if __name__ == "__main__":
    main()
